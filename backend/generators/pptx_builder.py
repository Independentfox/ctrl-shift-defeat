import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from models.context_object import StartupContextObject


BRAND_BLUE = RGBColor(0x1A, 0x73, 0xE8)
BRAND_DARK = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


def build_pitch_deck(slides_data: list[dict], context: StartupContextObject) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info.get("title", f"Slide {slide_info.get('slide_number', '')}")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BRAND_DARK

        # Content
        content = slide_info.get("content", "")
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for line in content.split("\n"):
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = BRAND_DARK
            p.space_after = Pt(8)

        # Source badge
        source = slide_info.get("source_agent", "")
        if source:
            badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(5), Inches(0.4))
            tf = badge_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Source: {source}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
